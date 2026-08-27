"""Génère le diaporama de soutenance de stage (Ascale).

Diapositives :
 1. Titre / Couverture
 2. Sommaire
 3. Présentation de l'entreprise
 4. Déroulement du stage — Équipe Agile / Daily meeting / Sprints
 5. Tâches confiées & Projet Ascale Marbre
 6. Architecture complète — Schéma en couches
 7. Architecture complète — Flux des requêtes
 8. Front-end
 9. Back-end
10. Chatbot LLM (OpenRouter)
11. Tests unitaires
12. Base de données & Schéma
13. API REST
14. Déploiement
15. Outils
16. Défis rencontrés
17. Conclusion & Démo
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Cm

# ── Palette ──────────────────────────────────────────
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
GOLD   = RGBColor(0xB8, 0x96, 0x4E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF7, 0xF6, 0xF2)
GREY   = RGBColor(0x88, 0x88, 0x88)
GREEN  = RGBColor(0x2E, 0xA0, 0x44)
BLUE   = RGBColor(0x1A, 0x6E, 0xC8)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank = prs.slide_layouts[6]   # totalement vide
    return prs.slides.add_slide(blank)


def fill_bg(slide, color):
    """Remplis le fond de la slide."""
    from pptx.oxml.ns import qn
    from lxml import etree
    sp = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, W, H
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.name = "_bg"
    # Met le rectangle en arrière-plan
    slide.shapes._spTree.remove(sp._element)
    slide.shapes._spTree.insert(2, sp._element)


def add_textbox(slide, text, x, y, w, h,
                font_size=20, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    sp = slide.shapes.add_shape(1, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill_color
    if line_color:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    return sp


def accent_bar(slide, y=Inches(6.8), color=GOLD):
    """Barre décorative dorée en bas."""
    add_rect(slide, 0, y, W, Inches(0.08), color)


def slide_header(slide, title, subtitle=None):
    """Bandeau titre + sous-titre."""
    add_rect(slide, 0, 0, W, Inches(1.4), DARK)
    add_textbox(slide, title, Inches(0.4), Inches(0.12), Inches(12), Inches(0.75),
                font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                    font_size=14, italic=True, color=GOLD, align=PP_ALIGN.LEFT)
    accent_bar(slide)


def bullet_block(slide, items, x, y, w, h, font_size=16, color=DARK, bullet="•"):
    """Bloc de bullets simple."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def code_block(slide, text, x, y, w, h):
    """Bloc monospace sur fond sombre."""
    add_rect(slide, x, y, w, h, DARK)
    txb = slide.shapes.add_textbox(
        x + Inches(0.1), y + Inches(0.1),
        w - Inches(0.2), h - Inches(0.2)
    )
    tf = txb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Courier New"
        run.font.size = Pt(10)
        run.font.color.rgb = GOLD


# ════════════════════════════════════════════════
prs = new_prs()


# ── 1. COUVERTURE ────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, DARK)
accent_bar(s, y=Inches(6.9))
# Logo / nom entreprise
add_textbox(s, "ASCALE", Inches(1), Inches(1.0), Inches(11), Inches(1.4),
            font_size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_textbox(s, "Importateur de matériaux de prestige — Marbre · Granit · Onyx · Travertin",
            Inches(1), Inches(2.3), Inches(11), Inches(0.5),
            font_size=16, italic=True, color=LIGHT, align=PP_ALIGN.CENTER)
add_rect(s, Inches(4.5), Inches(2.9), Inches(4.3), Inches(0.04), GOLD)
add_textbox(s, "RAPPORT DE STAGE — SOUTENANCE",
            Inches(1), Inches(3.1), Inches(11), Inches(0.6),
            font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s,
    "Salma [NOM]  ·  BUT Informatique 2ème année  ·  IUT de Villetaneuse\n"
    "Maître de stage : Nasser  ·  Tutrice : Pascale Hellegouarch  ·  Année 2025–2026",
    Inches(1), Inches(3.85), Inches(11), Inches(0.9),
    font_size=14, italic=True, color=GREY, align=PP_ALIGN.CENTER)


# ── 2. SOMMAIRE ──────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT)
slide_header(s, "Sommaire")

items_col1 = [
    "01 · Présentation de l'entreprise",
    "02 · Déroulement du stage",
    "03 · Tâches confiées & Projet Ascale",
    "04 · Architecture complète",
    "05 · Front-end",
    "06 · Back-end",
    "07 · Chatbot LLM (OpenRouter)",
    "08 · Système de paiement simulé",
]
items_col2 = [
    "09 · Tests unitaires",
    "10 · Base de données & Schéma",
    "11 · API REST",
    "12 · Déploiement",
    "13 · Outils",
    "14 · Défis rencontrés",
    "15 · Conclusion & Démo",
]
bullet_block(s, items_col1, Inches(0.5), Inches(1.6), Inches(6), Inches(5.5),
             font_size=17, color=DARK, bullet="▸")
bullet_block(s, items_col2, Inches(6.8), Inches(1.6), Inches(6), Inches(5.5),
             font_size=17, color=DARK, bullet="▸")


# ── 3. PRÉSENTATION ENTREPRISE ───────────────────────
s = blank_slide(prs)
fill_bg(s, WHITE)
slide_header(s, "Présentation de l'entreprise", "Ascale — spécialiste des pierres naturelles de prestige")

# 4 blocs catégories
cats = [
    ("Marbre",    "10 références\n250–1 380 MAD/m²"),
    ("Granit",    "4 références\n480–780 MAD/m²"),
    ("Onyx",      "3 références\n1 800–2 800 MAD/m²"),
    ("Travertin", "3 références\n280–350 MAD/m²"),
]
for i, (cat, desc) in enumerate(cats):
    x = Inches(0.4 + i * 3.2)
    add_rect(s, x, Inches(1.6), Inches(3.0), Inches(1.4), DARK)
    add_textbox(s, cat, x + Inches(0.1), Inches(1.65), Inches(2.8), Inches(0.6),
                font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(s, desc, x + Inches(0.1), Inches(2.2), Inches(2.8), Inches(0.7),
                font_size=13, color=LIGHT, align=PP_ALIGN.CENTER)

bullet_block(s, [
    "Clientèle : architectes d'intérieur, promoteurs immobiliers, artisans, particuliers",
    "Showroom à Casablanca (Zone Industrielle Oukacha) · Livraison sur tout le Maroc",
    "Avant le stage : gestion manuelle (Excel + WhatsApp) → digitalisation en cours",
    "Mission confiée : concevoir une application web de gestion commerciale complète",
], Inches(0.4), Inches(3.2), Inches(12.5), Inches(3.8),
   font_size=16, color=DARK)


# ── 4. DÉROULEMENT DU STAGE ──────────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT)
slide_header(s, "Déroulement du stage", "10 semaines · 5 sprints · équipe Agile Scrum")

# Sprints timeline
sprints = [
    ("S0", "Cadrage\n& Analyse"),
    ("S1", "Back-end\n& BDD"),
    ("S2", "Commandes\n& Emails"),
    ("S3", "Chatbot\n& Dashboard"),
    ("S4", "Front-end\n& Tests"),
]
for i, (num, desc) in enumerate(sprints):
    x = Inches(0.5 + i * 2.55)
    add_rect(s, x, Inches(1.55), Inches(2.3), Inches(0.6), DARK)
    add_textbox(s, f"Sprint {num}", x, Inches(1.6), Inches(2.3), Inches(0.5),
                font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_rect(s, x + Inches(0.9), Inches(2.2), Inches(0.5), Inches(0.5), GOLD)
    add_textbox(s, desc, x, Inches(2.8), Inches(2.3), Inches(0.7),
                font_size=11, color=DARK, align=PP_ALIGN.CENTER)

# Cérémonies Scrum
add_textbox(s, "Cérémonies Scrum", Inches(0.4), Inches(3.7), Inches(6), Inches(0.45),
            font_size=16, bold=True, color=DARK)
bullet_block(s, [
    "Daily meeting 9h00 (15 min) — avancement, blocages, plan du jour",
    "Sprint Planning — sélection user stories, estimation Fibonacci",
    "Sprint Review — démonstration au Product Owner (Nasser)",
    "Rétrospective — amélioration continue",
], Inches(0.4), Inches(4.1), Inches(6.2), Inches(3), font_size=14, color=DARK)

# Rôles
add_textbox(s, "Rôles", Inches(7), Inches(3.7), Inches(6), Inches(0.45),
            font_size=16, bold=True, color=DARK)
bullet_block(s, [
    "Product Owner : Nasser (maître de stage)",
    "Scrum Master : Nasser (double casquette)",
    "Développeuse : Salma (stagiaire)",
    "Outil : Azure DevOps (Boards + Repos + Pipelines)",
], Inches(7), Inches(4.1), Inches(5.9), Inches(3), font_size=14, color=DARK)


# ── 5. TÂCHES CONFIÉES & PROJET ──────────────────────
s = blank_slide(prs)
fill_bg(s, WHITE)
slide_header(s, "Tâches confiées & Projet Ascale Marbre",
             "Application web de gestion commerciale — Flask · SQLite · Python")

domaines = [
    ("Front-end",   ["Catalogue + filtres JS","Calculateur de surface","Chatbot flottant (AJAX)","Dashboard Chart.js"]),
    ("Back-end",    ["Routes Flask 3.0.3","Chatbot 10 intentions + LLM","Emails Flask-Mail","Webhook WhatsApp"]),
    ("Base données",["8 entités SQLAlchemy","Transactions atomiques","20 produits seed","Tests StaticPool"]),
    ("API & Dépl.", ["API JSON chatbot","Export CSV (BOM)","Déploiement Render.com","105 tests pytest"]),
]
for i, (dom, items) in enumerate(domaines):
    x = Inches(0.3 + i * 3.25)
    add_rect(s, x, Inches(1.5), Inches(3.1), Inches(0.5), GOLD)
    add_textbox(s, dom, x, Inches(1.52), Inches(3.1), Inches(0.45),
                font_size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    bullet_block(s, items, x + Inches(0.1), Inches(2.1), Inches(2.9), Inches(4.8),
                 font_size=13, color=DARK, bullet="✓")


# ── 6. ARCHITECTURE — SCHÉMA EN COUCHES ──────────────
s = blank_slide(prs)
fill_bg(s, DARK)
slide_header(s, "Architecture complète — Vue en couches (MVC)", "")

layers = [
    (WHITE,  "PRÉSENTATION (Vue)",   "Templates Jinja2 · CSS custom · JavaScript ES6\nbase_public.html · base.html · nos_materiaux · commander · chatbot · dashboard"),
    (GOLD,   "CONTRÔLEUR (Routes)",  "app.py · create_app(test_config)\nRoutes publiques + admin + API JSON + Export CSV + Webhook WhatsApp"),
    (WHITE,  "SERVICE (Modèle)",     "models.py · class Store  |  chatbot.py · repondre()  |  email_service.py\nTransactions atomiques · 10 intentions · LLM OpenRouter · Flask-Mail"),
    (GOLD,   "DONNÉES",              "database.py · 8 entités SQLAlchemy\nascale.db  (SQLite — fichier unique)"),
]
colors = [RGBColor(0x1A,0x6E,0xC8), RGBColor(0x2E,0xA0,0x44),
          RGBColor(0xB8,0x00,0x60), RGBColor(0x88,0x44,0x00)]
for i, (tc, title, desc) in enumerate(layers):
    y = Inches(1.45 + i * 1.35)
    add_rect(s, Inches(0.3), y, Inches(12.7), Inches(1.2),
             colors[i % 4])
    add_textbox(s, title, Inches(0.4), y + Inches(0.05), Inches(4), Inches(0.45),
                font_size=14, bold=True, color=WHITE)
    add_textbox(s, desc, Inches(4.5), y + Inches(0.05), Inches(8.4), Inches(1.1),
                font_size=12, color=LIGHT)


# ── 7. ARCHITECTURE — FLUX DES REQUÊTES ──────────────
s = blank_slide(prs)
fill_bg(s, WHITE)
slide_header(s, "Architecture complète — Flux des requêtes", "")

code_block(s,
"""Navigateur  ──HTTP GET/POST──▶  Flask (app.py)
                                     │
              ┌──────────────────────┼──────────────────────┐
              │                      │                      │
     Route publique         Route commande          Route chatbot
     /nos-materiaux         /commander              /chatbot/widget/envoyer
          │                      │                      │
     Store.list_produits()  Store.creer_commande()  chatbot.repondre()
          │                      │                    ├── mots-clés (10 intents)
          ▼                      ├──▶ email_service   └── _tenter_llm()
       SQLite                   │        └──▶ Gmail SMTP     └──▶ OpenRouter API
                                └──▶ SQLite

     Route admin              Webhook WhatsApp
     /gestion · /commandes    /webhook/whatsapp
          │                      │
     Dashboard (KPI)         whatsapp_config.py
     Chart.js                    └──▶ Meta Graph API v20.0""",
    Inches(0.3), Inches(1.5), Inches(12.7), Inches(5.5))


# ── 8. FRONT-END ─────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, WHITE)
slide_header(s, "Front-end", "Jinja2 · CSS custom · JavaScript ES6 · Chart.js")

bullet_block(s, [
    "Design system : palette ivoire / anthracite / or — sans framework CSS externe",
    "CSS Grid + Flexbox + media queries → responsive mobile-first",
    "Variables CSS (--gold, --dark, --light) pour cohérence admin ↔ public",
], Inches(0.4), Inches(1.5), Inches(12.5), Inches(1.5), font_size=15, color=DARK)

features = [
    ("Catalogue",         "Filtres JS (catégorie, dispo, prix)\nSans rechargement de page"),
    ("Calculateur",       "Longueur × largeur × taux chutes\nInjection auto dans le formulaire"),
    ("Chatbot widget",    "AJAX → /chatbot/widget/envoyer\nHistorique conversation flottant"),
    ("Dashboard",         "KPI cards + Chart.js\nTop 5 produits · Commandes par canal"),
]
for i, (feat, desc) in enumerate(features):
    x = Inches(0.3 + i * 3.25)
    add_rect(s, x, Inches(3.1), Inches(3.1), Inches(0.5), DARK)
    add_textbox(s, feat, x, Inches(3.12), Inches(3.1), Inches(0.45),
                font_size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(s, desc, x + Inches(0.1), Inches(3.7), Inches(2.9), Inches(1.2),
                font_size=13, color=DARK)


# ── 9. BACK-END ──────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT)
slide_header(s, "Back-end", "Flask 3.0.3 · Python 3.11 · SQLAlchemy · Flask-Mail · Werkzeug")

add_textbox(s, "app.py — Routes", Inches(0.4), Inches(1.5), Inches(5.5), Inches(0.4),
            font_size=15, bold=True, color=DARK)
bullet_block(s, [
    "create_app(test_config) — injection config tests",
    "Validation serveur (email, quantités, stocks)",
    "Flash messages, sessions signées (Werkzeug)",
    "Export CSV avec BOM UTF-8 (Excel)",
], Inches(0.4), Inches(1.95), Inches(5.5), Inches(2.5), font_size=14, color=DARK)

add_textbox(s, "models.py — Store", Inches(0.4), Inches(4.5), Inches(5.5), Inches(0.4),
            font_size=15, bold=True, color=DARK)
bullet_block(s, [
    "Façade : clients, produits, commandes, créneaux",
    "StockInsuffisantError, CreneauIndisponibleError",
    "Transactions atomiques (flush → commit)",
    "seed() : 7 clients, 20 produits, 7 commandes",
], Inches(0.4), Inches(4.95), Inches(5.5), Inches(2.2), font_size=14, color=DARK)

add_textbox(s, "email_service.py", Inches(7), Inches(1.5), Inches(5.7), Inches(0.4),
            font_size=15, bold=True, color=DARK)
bullet_block(s, [
    "Flask-Mail → Gmail SMTP (TLS)",
    "Confirmation commande & réservation",
    "Auto-réponse contact + pré-réponse chatbot",
    "MAIL_SUPPRESS_SEND=True en tests",
], Inches(7), Inches(1.95), Inches(5.7), Inches(2.5), font_size=14, color=DARK)

add_textbox(s, "whatsapp_config.py", Inches(7), Inches(4.5), Inches(5.7), Inches(0.4),
            font_size=15, bold=True, color=DARK)
bullet_block(s, [
    "Meta Graph API v20.0",
    "Webhook vérification (challenge/response)",
    "Parsing payloads JSON entrants",
    "Prêt — en attente credentials Meta",
], Inches(7), Inches(4.95), Inches(5.7), Inches(2.2), font_size=14, color=DARK)


# ── 10. CHATBOT LLM ──────────────────────────────────
s = blank_slide(prs)
fill_bg(s, DARK)
slide_header(s, "Chatbot LLM — OpenRouter", "Moteur mots-clés + fallback LLM Mistral 7B")

add_textbox(s, "Chaîne de 10 intentions (ordre de priorité)",
            Inches(0.4), Inches(1.45), Inches(6), Inches(0.45),
            font_size=15, bold=True, color=GOLD)
bullet_block(s, [
    "1. Salutation — bonjour, salam, hello…",
    "2. Politesse — merci, au revoir…",
    "3. Action — commander, réserver, showroom…",
    "4. Devis m² — regex surface + produit",
    "5. Statut commande — numéro ou nom client",
    "6. Budget — montant MAD ou qualitatif",
    "7. Recommandation — cuisine, salle de bain…",
    "8. FAQ — horaires, adresse, livraison…",
    "9. Info produit — score mots ≥ 2",
    "10. Catégorie — marbre, granit, onyx…",
    "→  Fallback LLM (si OPENROUTER_API_KEY)",
], Inches(0.4), Inches(1.95), Inches(5.8), Inches(5.2), font_size=13, color=LIGHT)

add_textbox(s, "LLM via OpenRouter API",
            Inches(7), Inches(1.45), Inches(6), Inches(0.45),
            font_size=15, bold=True, color=GOLD)
bullet_block(s, [
    "Modèle : mistralai/mistral-7b-instruct:free",
    "Appel HTTPS via urllib.request (stdlib)",
    "System prompt contraint au domaine matériaux",
    "Catalogue temps réel injecté dans le prompt",
    "Timeout 10s — fallback message d'aide si erreur",
    "OPENROUTER_API_KEY vide = désactivé (sécurité)",
], Inches(7), Inches(1.95), Inches(5.9), Inches(3), font_size=14, color=LIGHT)

code_block(s,
"""# Fallback LLM (chatbot.py)
def _tenter_llm(store, message):
    if not OPENROUTER_API_KEY:
        return None
    # Catalogue temps réel
    catalogue = "\\n".join(
      f"- {p.nom} : {p.prix:.0f} MAD/m²"
      for p in store.list_produits()
    )
    payload = {"model": OPENROUTER_MODEL,
               "messages": [
                 {"role":"system",
                  "content": SYSTEM_PROMPT.format(
                    catalogue=catalogue)},
                 {"role":"user",
                  "content": message}],
               "max_tokens": 350}
    # POST vers OpenRouter…""",
    Inches(7), Inches(4.85), Inches(6.1), Inches(2.3))


# ── 11. SYSTÈME DE PAIEMENT SIMULÉ ──────────────────
s = blank_slide(prs)
fill_bg(s, WHITE)
slide_header(s, "Système de paiement simulé",
             "Tunnel 3 étapes · Nominatim · fpdf2 · Flask-Mail")

# Tunnel 3 étapes
etapes = [
    ("1", "Produits",      "Sélection produits\nInfos contact"),
    ("2", "Paiement",      "Adresse livraison\n+ mode de règlement"),
    ("3", "Confirmation",  "Reçu immédiat\n+ facture PDF par email"),
]
for i, (num, titre, desc) in enumerate(etapes):
    x = Inches(0.5 + i * 4.2)
    add_rect(s, x, Inches(1.5), Inches(3.8), Inches(0.6), DARK)
    add_textbox(s, f"Étape {num} — {titre}", x, Inches(1.55), Inches(3.8), Inches(0.5),
                font_size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(s, desc, x + Inches(0.1), Inches(2.25), Inches(3.6), Inches(1.0),
                font_size=13, color=DARK, align=PP_ALIGN.CENTER)

# Deux colonnes
add_textbox(s, "Validation adresse — Nominatim / OSM",
            Inches(0.4), Inches(3.4), Inches(6.0), Inches(0.4),
            font_size=15, bold=True, color=DARK)
bullet_block(s, [
    "Autocomplétion JS (debounce 400 ms) — suggestions en temps réel",
    "Validation côté serveur via urllib.request (stdlib, 0 dépendance)",
    "Fail-open : jamais bloquant si Nominatim indisponible",
    "OpenStreetMap — sans clé API, sans coût",
], Inches(0.4), Inches(3.85), Inches(6.0), Inches(2.8), font_size=14, color=DARK)

add_textbox(s, "Modes de paiement & facturation",
            Inches(7), Inches(3.4), Inches(5.9), Inches(0.4),
            font_size=15, bold=True, color=DARK)
bullet_block(s, [
    "💳 Carte bancaire — champs visuels (aucun traitement réel)",
    "🏦 Virement bancaire — RIB Ascale affiché",
    "💰 Espèces au showroom — confirmation directe",
    "Facture pro forma PDF générée en mémoire (fpdf2 2.8.7)",
    "Email avec PDF joint via Flask-Mail (SMTP Gmail)",
], Inches(7), Inches(3.85), Inches(5.9), Inches(2.8), font_size=14, color=DARK)


# ── 12. TESTS UNITAIRES ──────────────────────────────
s = blank_slide(prs)
fill_bg(s, WHITE)
slide_header(s, "Tests unitaires", "105 tests · pytest + pytest-flask · SQLite in-memory")

# Stats visuelles
stats = [("42", "tests\nmodèles"), ("27", "tests\nchatbot"), ("36", "tests\nroutes"), ("100%", "passent")]
for i, (num, label) in enumerate(stats):
    x = Inches(0.5 + i * 3.2)
    add_rect(s, x, Inches(1.5), Inches(2.8), Inches(1.5),
             DARK if i < 3 else GOLD)
    add_textbox(s, num, x, Inches(1.6), Inches(2.8), Inches(0.9),
                font_size=44, bold=True,
                color=GOLD if i < 3 else DARK, align=PP_ALIGN.CENTER)
    add_textbox(s, label, x, Inches(2.45), Inches(2.8), Inches(0.45),
                font_size=13, color=WHITE if i < 3 else DARK, align=PP_ALIGN.CENTER)

bullet_block(s, [
    "Isolation : StaticPool — même connexion in-memory pour tout le test → pas d'effet de bord",
    "scope='function' — base SQLite vierge fraîche pour chaque test",
    "Mocks : patch('chatbot.OPENROUTER_API_KEY', '') — LLM désactivé en test",
    "test_models.py — atomicité commandes, créneaux, stock_faible, rupture, list_tous_creneaux",
    "test_chatbot.py — 10 intentions, 4 faux positifs, 3 scénarios LLM mock",
    "test_routes.py — GET/POST, codes HTTP, réponse JSON, export CSV, webhook WhatsApp",
], Inches(0.4), Inches(3.15), Inches(12.5), Inches(4), font_size=14, color=DARK)


# ── 12. BASE DE DONNÉES & SCHÉMA ─────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT)
slide_header(s, "Base de données & Schéma", "SQLite · Flask-SQLAlchemy · 8 entités · 20 produits seed")

code_block(s,
"""CLIENT ─(1,N)──▶ COMMANDE ─(1,N)──▶ LIGNE_COMMANDE ─(N,1)──▶ PRODUIT ─(N,1)──▶ CATEGORIE
                                                                      │
CLIENT ─(0,N)──▶ CRENEAU ─(0,1)──▶ RESERVATION                  stock/prix

UTILISATEUR ── accède à ── [Espace Admin]

Entités : CLIENT · PRODUIT · CATEGORIE · COMMANDE
          LIGNE_COMMANDE · CRENEAU · RESERVATION · UTILISATEUR""",
    Inches(0.3), Inches(1.5), Inches(12.7), Inches(2.0))

bullet_block(s, [
    "db.session.flush() — obtenir l'ID avant commit (transactions atomiques multi-lignes)",
    "StockInsuffisantError — vérification TOUS les stocks avant toute décrémentation",
    "20 produits seed (10 marbres · 4 granits · 3 onyx · 3 travertins) · 7 clients · 7 commandes",
    "SQLAlchemy ORM — migration PostgreSQL facilitée (seul SQLALCHEMY_DATABASE_URI change)",
], Inches(0.4), Inches(3.7), Inches(12.5), Inches(3.5), font_size=15, color=DARK)


# ── 13. API REST ──────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, WHITE)
slide_header(s, "API REST", "Routes publiques · admin · JSON · CSV · Webhook")

endpoints = [
    ("GET  /nos-materiaux",        "HTML", "Catalogue produits + filtres"),
    ("POST /commander",            "HTML", "Création commande publique → redirect"),
    ("POST /chatbot/widget/envoyer","JSON","API chatbot (widget flottant AJAX)"),
    ("GET  /admin/commandes/export","CSV", "Export commandes (UTF-8 BOM)"),
    ("GET  /admin/clients/export",  "CSV", "Export clients (UTF-8 BOM)"),
    ("POST /commandes/<id>/statut", "HTML","Mise à jour statut dropdown inline"),
    ("GET  /webhook/whatsapp",     "TEXT", "Vérification webhook Meta (challenge)"),
    ("POST /webhook/whatsapp",     "JSON", "Réception messages WhatsApp Business"),
]
# ─ Table via textboxes ─
col_w = [4.3, 1.5, 7]
y_start = Inches(1.5)
row_h   = Inches(0.58)
# Header
for ci, (hdr, cw) in enumerate(zip(["Endpoint", "Format", "Description"], col_w)):
    cx = Inches(0.3 + sum(col_w[:ci]))
    add_rect(s, cx, y_start, Inches(cw), row_h, DARK)
    add_textbox(s, hdr, cx + Inches(0.05), y_start + Inches(0.08), Inches(cw - 0.1), row_h,
                font_size=12, bold=True, color=WHITE)
for ri, row in enumerate(endpoints):
    bg = LIGHT if ri % 2 == 0 else WHITE
    for ci, (cell, cw) in enumerate(zip(row, col_w)):
        cx = Inches(0.3 + sum(col_w[:ci]))
        cy = y_start + row_h * (ri + 1)
        add_rect(s, cx, cy, Inches(cw), row_h, bg)
        fc = GOLD if ci == 0 else DARK
        add_textbox(s, cell, cx + Inches(0.05), cy + Inches(0.08), Inches(cw - 0.1), row_h,
                    font_size=11, color=fc)


# ── 14. DÉPLOIEMENT ───────────────────────────────────
s = blank_slide(prs)
fill_bg(s, DARK)
slide_header(s, "Déploiement prévu", "Render.com · Gunicorn · GitHub CI · Persistent Disk")

add_textbox(s, "Solution retenue : Render.com (Free tier)",
            Inches(0.4), Inches(1.45), Inches(12), Inches(0.45),
            font_size=18, bold=True, color=GOLD)

code_block(s,
"""GitHub (main) ──git push──▶ Render.com CI
                                    │
                           Build : pip install -r requirements.txt
                                   (+ gunicorn)
                                    │
                           Start  : gunicorn app:app
                                    │
                    ┌───────────────┼───────────────────┐
                    │               │                   │
              Env vars          Persistent Disk     HTTPS auto
              SECRET_KEY        /data/ascale.db     ascale.onrender.com
              MAIL_*            (5 Go gratuit)
              OPENROUTER_API_KEY""",
    Inches(0.3), Inches(2.0), Inches(12.7), Inches(3.2))

bullet_block(s, [
    "Avantage : déploiement par git push, sans configuration nginx/gunicorn manuelle",
    "Inconvénient tier free : mise en veille après 15 min → upgrade $7/mois si besoin",
    "Alternative : PythonAnywhere (facile Flask, console SSH, HTTPS inclus)",
], Inches(0.4), Inches(5.4), Inches(12.5), Inches(2), font_size=14, color=LIGHT)


# ── 15. OUTILS ───────────────────────────────────────
s = blank_slide(prs)
fill_bg(s, LIGHT)
slide_header(s, "Outils & Stack technique", "")

outils = [
    ("Back-end",    "Python 3.11 · Flask 3.0.3\nSQLAlchemy 3.1 · Werkzeug\npython-dotenv · Gunicorn"),
    ("Front-end",   "Jinja2 · Vanilla CSS ES6\nChart.js 4.4 · Feather Icons\nCSS Grid + Flexbox"),
    ("Tests",       "pytest 8.3 · pytest-flask\nStaticPool · unittest.mock\n105 tests unitaires"),
    ("DevOps",      "Git · Azure DevOps\nBoards Scrum · PR Reviews\nCI Pipelines · Render.com"),
    ("Externe API", "OpenRouter (LLM)\nMeta Graph API v20.0\nFlask-Mail · Gmail SMTP"),
    ("IDE",         "VS Code\nPylance · GitLens\nSQLite Viewer · Chrome DevTools"),
]
for i, (cat, desc) in enumerate(outils):
    row, col = divmod(i, 3)
    x = Inches(0.3 + col * 4.35)
    y = Inches(1.5 + row * 2.6)
    add_rect(s, x, y, Inches(4.1), Inches(0.48), DARK)
    add_textbox(s, cat, x + Inches(0.1), y + Inches(0.04), Inches(3.9), Inches(0.4),
                font_size=14, bold=True, color=GOLD)
    add_textbox(s, desc, x + Inches(0.1), y + Inches(0.55), Inches(3.9), Inches(1.9),
                font_size=12, color=DARK)


# ── 16. DÉFIS RENCONTRÉS ─────────────────────────────
s = blank_slide(prs)
fill_bg(s, WHITE)
slide_header(s, "Défis rencontrés", "")

defis = [
    ("Migration BDD",          "Garder l'API publique identique\n→ Pattern Façade + flush()"),
    ("Chatbot faux positifs",  "10 intentions ordonnées\n→ Score mots entiers + seuil min"),
    ("LLM fallback",           "OpenRouter gratuit\n→ System prompt contraint au domaine"),
    ("Atomicité stock",        "Race condition commandes\n→ Validation TOUS avant décrémentation"),
    ("Tests isolés",           "DB partagée entre tests\n→ StaticPool + scope='function'"),
    ("WhatsApp Business",      "Credentials Meta non obtenus\n→ Flask-Mail en fallback immédiat"),
]
for i, (titre, sol) in enumerate(defis):
    row, col = divmod(i, 3)
    x = Inches(0.3 + col * 4.35)
    y = Inches(1.5 + row * 2.8)
    add_rect(s, x, y, Inches(4.1), Inches(0.48), GOLD)
    add_textbox(s, titre, x + Inches(0.1), y + Inches(0.04), Inches(3.9), Inches(0.4),
                font_size=14, bold=True, color=DARK)
    add_textbox(s, sol, x + Inches(0.1), y + Inches(0.55), Inches(3.9), Inches(2.1),
                font_size=13, color=DARK)


# ── 17. CONCLUSION & DÉMO ────────────────────────────
s = blank_slide(prs)
fill_bg(s, DARK)
accent_bar(s)
add_rect(s, 0, 0, W, Inches(0.08), GOLD)

add_textbox(s, "Conclusion & Démo",
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

bullet_block(s, [
    "Application complète livrée : catalogue · commandes · chatbot LLM · dashboard · tests",
    "Tunnel paiement simulé (3 étapes) · validation Nominatim · facture PDF · email pièce jointe",
    "Chatbot mode hybride : Spécialisé (10 intentions) ↔ IA Libre (OpenRouter Mistral 7B)",
    "105 tests unitaires passent — couverture back-end complète (modèles, chatbot, routes)",
    "20 produits (4 familles) · 7 clients · architecture MVC propre et extensible",
    "Déploiement Render.com documenté — prêt à mettre en ligne",
], Inches(0.5), Inches(1.1), Inches(12.3), Inches(4.0), font_size=15, color=LIGHT)

add_rect(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.06), GOLD)

add_textbox(s, "DÉMONSTRATION LIVE",
            Inches(0.5), Inches(4.9), Inches(12.3), Inches(0.7),
            font_size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

demo_items = [
    "① Catalogue public + filtres  |  ② Calculateur + tunnel paiement + facture PDF",
    "③ Chatbot mode Spécialisé / IA Libre  |  ④ Dashboard admin (KPI + graphiques)",
    "⑤ Export CSV  |  ⑥ Réservation showroom",
]
bullet_block(s, demo_items, Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.6),
             font_size=15, color=WHITE, bullet="")

add_textbox(s, "Merci de votre attention ! 🏛️",
            Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
            font_size=14, italic=True, color=GREY, align=PP_ALIGN.CENTER)


# ─── SAUVEGARDE ──────────────────────────────────────
out = r"C:\Users\EMI\Documents\salma stage\diaporama_stage_ascale.pptx"
prs.save(out)
print(f"OK — {out}")
